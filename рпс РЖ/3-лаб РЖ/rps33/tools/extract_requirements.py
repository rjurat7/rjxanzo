from __future__ import annotations

from pathlib import Path


def extract_docx(docx_path: Path) -> list[str]:
    from docx import Document

    doc = Document(str(docx_path))
    lines: list[str] = []

    for p in doc.paragraphs:
        t = (p.text or "").strip()
        if t:
            lines.append(t)

    for table in doc.tables:
        for row in table.rows:
            cells = []
            for c in row.cells:
                t = (c.text or "").replace("\n", " ").strip()
                if t:
                    cells.append(t)
            if cells:
                lines.append(" | ".join(cells))

    return lines


def extract_pptx(pptx_path: Path) -> list[str]:
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    lines: list[str] = []

    for i, slide in enumerate(prs.slides, start=1):
        lines.append(f"===== SLIDE {i} =====")
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            txt = (shape.text or "").strip()
            if not txt:
                continue
            for ln in txt.splitlines():
                ln = ln.strip()
                if ln:
                    lines.append(ln)

    return lines


def main() -> int:
    out_dir = Path(r"d:\rps-main\rps33\_requirements_extract")
    out_dir.mkdir(parents=True, exist_ok=True)

    docx_path = Path(r"c:\Users\DILSHOD\Downloads\Лабораторная_работа_3 (1).docx")
    (out_dir / "lab3_docx.txt").write_text("\n".join(extract_docx(docx_path)), encoding="utf-8")

    pptx_db = Path(r"c:\Users\DILSHOD\Downloads\Проектирование баз данных.pptx")
    (out_dir / "db_design_pptx.txt").write_text("\n".join(extract_pptx(pptx_db)), encoding="utf-8")

    pptx_ui = Path(r"c:\Users\DILSHOD\Downloads\Проектирование пользовательского интерфейса (1).pptx")
    (out_dir / "ui_design_pptx.txt").write_text("\n".join(extract_pptx(pptx_ui)), encoding="utf-8")

    print(f"Extracted to: {out_dir}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

