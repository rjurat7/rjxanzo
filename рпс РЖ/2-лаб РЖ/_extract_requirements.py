from pathlib import Path

out_dir = Path(r'd:\rps2\_requirements_extracted')
out_dir.mkdir(parents=True, exist_ok=True)

# DOCX
from docx import Document

docx_path = Path(r'c:\Users\DILSHOD\Downloads\Лабораторная работа 2 (3).docx')
doc = Document(str(docx_path))
doc_text = []
for p in doc.paragraphs:
    t = (p.text or '').strip()
    if t:
        doc_text.append(t)

for table in doc.tables:
    for row in table.rows:
        cells = [(c.text or '').strip() for c in row.cells]
        line = ' | '.join([c for c in cells if c])
        if line:
            doc_text.append(line)

(out_dir / 'lab2_docx.txt').write_text('\n'.join(doc_text), encoding='utf-8')

# PPTX
from pptx import Presentation

pptx_path = Path(r'c:\Users\DILSHOD\Downloads\Расчет сложности алгоритма (1).pptx')
prs = Presentation(str(pptx_path))
ppt_lines = []
for i, slide in enumerate(prs.slides, start=1):
    slide_lines = []
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            t = (shape.text or '').strip()
            if t:
                slide_lines.append(t)
    if slide_lines:
        ppt_lines.append(f'--- SLIDE {i} ---')
        ppt_lines.extend(slide_lines)

(out_dir / 'algo_complexity_pptx.txt').write_text('\n'.join(ppt_lines), encoding='utf-8')

print('OK', out_dir)
